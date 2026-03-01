"""
Generate product catalogue for Campaign Analytics Lab (public demo).
Run: python make_catalogue.py
Output: Campaign_Analytics_Lab_Catalogue.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

NAVY   = RGBColor(0x1a, 0x1a, 0x6b)
ORANGE = RGBColor(0xF4, 0x79, 0x20)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF8, 0xF8, 0xF8)
GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
DARK   = RGBColor(0x11, 0x11, 0x11)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
TEAL   = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

def _slide():
    return prs.slides.add_slide(BLANK)

def _rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def _txt(slide, text, left, top, width, height,
         size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
         italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def _accent_bar(slide, top, color=ORANGE, height=Inches(0.04)):
    _rect(slide, 0, top, W, height, color)

def _ph(slide, left, top, width, height, label):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF8)
    box.line.color.rgb = NAVY; box.line.width = Pt(1.2)
    tf = box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"📸  {label}"
    r.font.size = Pt(10); r.font.color.rgb = NAVY; r.font.italic = True

def _card(slide, left, top, width, height, icon, title, body,
          bg=WHITE, title_color=DARK, body_color=RGBColor(0x44,0x44,0x44)):
    _rect(slide, left, top, width, height, bg)
    border = slide.shapes.add_shape(1, left, top, Inches(0.04), height)
    border.fill.solid(); border.fill.fore_color.rgb = ORANGE; border.line.fill.background()
    _txt(slide, icon, left+Inches(0.12), top+Inches(0.08), Inches(0.4), Inches(0.4), size=18, color=title_color)
    _txt(slide, title, left+Inches(0.55), top+Inches(0.10), width-Inches(0.65), Inches(0.35), size=11, bold=True, color=title_color)
    _txt(slide, body, left+Inches(0.12), top+Inches(0.48), width-Inches(0.22), height-Inches(0.6), size=9, color=body_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, NAVY)
_rect(sl, 0, H-Inches(0.12), W, Inches(0.12), ORANGE)
_rect(sl, W*0.55, 0, W*0.45, H, RGBColor(0x0e,0x0e,0x55))

_txt(sl, "PRODUCT OVERVIEW", Inches(0.7), Inches(0.8), Inches(5), Inches(0.4),
     size=9, bold=True, color=ORANGE)
_txt(sl, "Campaign Analytics Lab", Inches(0.7), Inches(1.3), Inches(6), Inches(1.2),
     size=36, bold=True, color=WHITE)
_txt(sl, "Decision Intelligence for Marketing Teams", Inches(0.7), Inches(2.55),
     Inches(5.8), Inches(0.5), size=16, color=GRAY)
_txt(sl, (
    "Model campaign scenarios, isolate demand shocks,\n"
    "and project outcomes — all in a browser, no code required."
), Inches(0.7), Inches(3.2), Inches(5.5), Inches(1.2),
    size=12, color=RGBColor(0xCC,0xCC,0xCC))

_txt(sl, "🌐  Try it free:  https://campaign-analytics-lab.streamlit.app",
     Inches(0.7), Inches(4.7), Inches(5.5), Inches(0.4), size=10, bold=True, color=ORANGE)
_txt(sl, "Demo login:  demo  /  demo2026",
     Inches(0.7), Inches(5.15), Inches(5.5), Inches(0.35), size=9, color=GRAY)

_txt(sl, "📸  INSERT: App screenshot (public URL — login or dashboard)",
     Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.2),
     size=10, color=RGBColor(0x88,0x88,0xBB), italic=True, align=PP_ALIGN.CENTER)

_txt(sl, f"{datetime.date.today().year}  ·  Open-source demo",
     Inches(0.7), Inches(6.9), Inches(4), Inches(0.4), size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IT DOES
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, Inches(1.3), NAVY)
_rect(sl, 0, Inches(1.3), W, H-Inches(1.3), LIGHT)
_accent_bar(sl, Inches(1.3))

_txt(sl, "What It Does", Inches(0.7), Inches(0.22), Inches(8), Inches(0.6),
     size=26, bold=True, color=WHITE)
_txt(sl, "Four core capabilities — available out of the box",
     Inches(0.7), Inches(0.78), Inches(8), Inches(0.4), size=11, color=GRAY)

caps = [
    ("📊", "Demand Forecasting",
     "Seasonal DNA engine reconstructs your organic demand baseline from historical data. "
     "Enter trial-period actuals to calibrate and generate a 12-month projection."),
    ("⚡", "Campaign Simulation",
     "Model campaigns with four response shapes (Front-Loaded, Linear Fade, Step, Delayed Peak). "
     "See immediate impact on clicks, conversions, and revenue."),
    ("🧹", "De-Shock Analysis",
     "Isolate artificial spikes from historical or forecast data, extract them as "
     "signatures, and re-inject on any future date."),
    ("📋", "Attribution Engine",
     "Waterfall breakdown of which events contribute how much to your revenue target — "
     "in absolute volume and percentage of gap."),
]
for i, (ic, tt, dd) in enumerate(caps):
    col = i % 2; row = i // 2
    _card(sl,
          left=Inches(0.5+col*6.3),
          top=Inches(1.7+row*2.55),
          width=Inches(6.0), height=Inches(2.25),
          icon=ic, title=tt, body=dd)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DASHBOARD SCREENSHOT
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, W, Inches(0.6), NAVY)
_txt(sl, "📊  Dashboard", Inches(0.5), Inches(0.1),
     Inches(8), Inches(0.4), size=14, bold=True, color=WHITE)

_ph(sl, Inches(0.3), Inches(0.75), Inches(8.5), Inches(5.5),
    "INSERT REAL SCREENSHOT from https://campaign-analytics-lab.streamlit.app — Dashboard view")

for i, b in enumerate([
    "Baseline (organic) vs Simulation (events) projection",
    "Monthly / Weekly / Daily granularity toggle",
    "Multi-entity (brand) blended DNA",
    "Goal Tracker with growth % and volume driver selector",
]):
    _txt(sl, f"✓  {b}", Inches(9.1), Inches(1.0+i*0.85),
         Inches(3.9), Inches(0.65), size=10, color=DARK)

_rect(sl, Inches(9.0), Inches(4.5), Inches(4.1), Inches(1.8), NAVY)
_txt(sl, "Try it now:", Inches(9.1), Inches(4.62), Inches(3.9), Inches(0.4),
     size=9, bold=True, color=ORANGE)
_txt(sl, "https://campaign-analytics-lab\n  .streamlit.app",
     Inches(9.1), Inches(5.05), Inches(3.9), Inches(1.1), size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SIMULATION LAB SCREENSHOT
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, LIGHT)
_rect(sl, 0, 0, W, Inches(0.6), NAVY)
_txt(sl, "⚡  Simulation Lab", Inches(0.5), Inches(0.1),
     Inches(8), Inches(0.4), size=14, bold=True, color=WHITE)

_ph(sl, Inches(0.3), Inches(0.75), Inches(8.5), Inches(5.5),
    "INSERT REAL SCREENSHOT — Simulation Lab (Events or De-Shock tab)")

tabs = [
    ("🖱️", "Visual DNA Drag",   "Reshape seasonal demand indices interactively"),
    ("🚀", "Event Injection",   "Campaigns, DNA swaps, custom drag events"),
    ("🧹", "De-Shock Tool",    "Historical & forecast spike isolation + compressor"),
    ("📋", "Audit & Attribution", "Waterfall event attribution with gap %"),
]
for i, (ic, tt, dd) in enumerate(tabs):
    top = Inches(0.85 + i * 1.5)
    _rect(sl, Inches(9.1), top, Inches(4.0), Inches(1.3), WHITE)
    border = sl.shapes.add_shape(1, Inches(9.1), top, Inches(0.04), Inches(1.3))
    border.fill.solid(); border.fill.fore_color.rgb = ORANGE; border.line.fill.background()
    _txt(sl, ic, Inches(9.22), top+Inches(0.1), Inches(0.4), Inches(0.4), size=14, color=DARK)
    _txt(sl, tt, Inches(9.65), top+Inches(0.1), Inches(3.3), Inches(0.35), size=10, bold=True, color=DARK)
    _txt(sl, dd, Inches(9.22), top+Inches(0.52), Inches(3.8), Inches(0.65), size=8.5, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY FEATURES GRID
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, Inches(1.3), NAVY)
_rect(sl, 0, Inches(1.3), W, H-Inches(1.3), WHITE)
_accent_bar(sl, Inches(1.3))

_txt(sl, "Key Features", Inches(0.7), Inches(0.22), Inches(10), Inches(0.6),
     size=26, bold=True, color=WHITE)
_txt(sl, "Everything included — one deployment, zero add-ons",
     Inches(0.7), Inches(0.78), Inches(9), Inches(0.4), size=11, color=GRAY)

feats = [
    ("🧬", "Seasonal DNA Engine",      "Weighted multi-year seasonality for clicks, CR, and AOV"),
    ("🎛️", "Spike Compressor",         "Scale any spike ±% with Gaussian edge taper"),
    ("🔬", "Brand Forge",              "Synthesise new entities from DNA inheritance + volume targets"),
    ("📁", "Excel Export",             "One-click formatted strategy report download"),
    ("🌍", "Multi-entity support",     "Compare and blend unlimited entity DNA profiles"),
    ("⚙️", "Configurable defaults",    "Per-entity campaign coefficients with global override"),
    ("🔐", "Auth gate",                "Username + password session, full activity log"),
    ("📖", "Built-in documentation",   "Full model explanation and formula guide in-app"),
    ("📊", "Goal Tracker",             "Set revenue / orders / clicks targets with growth scenarios"),
]
for i, (ic, tt, dd) in enumerate(feats):
    col = i % 3; row = i // 3
    _card(sl,
          left=Inches(0.45+col*4.28),
          top=Inches(1.52+row*1.95),
          width=Inches(3.95), height=Inches(1.75),
          icon=ic, title=tt, body=dd)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW TO GET STARTED
# ══════════════════════════════════════════════════════════════════════════════
sl = _slide()
_rect(sl, 0, 0, W, H, NAVY)
_rect(sl, 0, H-Inches(0.12), W, Inches(0.12), ORANGE)

_txt(sl, "Get Started in 3 Steps", Inches(0.7), Inches(0.5),
     Inches(12), Inches(0.7), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

steps = [
    ("1", "Open the demo", "Visit https://campaign-analytics-lab.streamlit.app\nLogin: demo / demo2026"),
    ("2", "Enter your trial data", "Set trial period dates + actual clicks, orders, and revenue in the sidebar"),
    ("3", "Simulate & Export", "Inject campaigns, run the de-shock tool, download your strategy report"),
]
for i, (num, title, body) in enumerate(steps):
    left = Inches(0.5 + i * 4.27)
    _rect(sl, left, Inches(1.6), Inches(3.95), Inches(4.2), RGBColor(0x14,0x14,0x55))
    circle = sl.shapes.add_shape(9, left+Inches(1.55), Inches(1.8), Inches(0.85), Inches(0.85))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE; circle.line.fill.background()
    _txt(sl, num, left+Inches(1.55), Inches(1.8), Inches(0.85), Inches(0.85),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, title, left+Inches(0.2), Inches(2.85), Inches(3.55), Inches(0.5),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, body, left+Inches(0.2), Inches(3.45), Inches(3.55), Inches(2.1),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

_rect(sl, Inches(3.5), Inches(6.15), Inches(6.33), Inches(0.85),
      RGBColor(0xF4,0x79,0x20))
_txt(sl, "🔗  github.com/Parsa-Hajian/campaign-analytics-lab  ·  Open source",
     Inches(3.5), Inches(6.22), Inches(6.33), Inches(0.6),
     size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "Campaign_Analytics_Lab_Catalogue.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
